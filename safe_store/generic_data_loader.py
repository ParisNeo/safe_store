from pathlib import Path
from typing import List


class PackageManager:
    @staticmethod
    def install_package(package_name):
        import subprocess
        import sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])
 
class GenericDataLoader:
    @staticmethod
    def read_file(file_path: Path|str) -> str:
        """
        Read a file and return its content as a string.

        Args:
            file_path (Path): The path to the file.

        Returns:
            str: The content of the file.
        
        Raises:
            ValueError: If the file type is unknown.
        """
        file_path = Path(file_path)
        
        if file_path.suffix.lower() ==".pdf":
            return GenericDataLoader.read_pdf_file(file_path)
        elif file_path.suffix.lower() == ".docx":
            return GenericDataLoader.read_docx_file(file_path)
        elif file_path.suffix.lower() == ".json":
            return GenericDataLoader.read_json_file(file_path)
        elif file_path.suffix.lower() == ".html":
            return GenericDataLoader.read_html_file(file_path)
        elif file_path.suffix.lower() == ".pptx":
            return GenericDataLoader.read_pptx_file(file_path)
        if file_path.suffix.lower() in [".pcap"]:
            return GenericDataLoader.read_pcap_file(file_path)
        if file_path.suffix.lower() in [".txt", ".rtf", ".md", ".log", ".csv", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat"]:
            return GenericDataLoader.read_text_file(file_path)
        else:
            raise ValueError("Unknown file type")
        
    @staticmethod
    def get_supported_file_types() -> List[str]:
        """
        Get the list of supported file types.

        Returns:
            List[str]: The list of supported file types.
        """
        return ["pdf", "txt", "docx", "json", "css", "css", "html", "pptx",".txt", ".md", ".log", ".cpp", ".java", ".js", ".py", ".rb", ".sh", ".sql", ".css", ".html", ".php", ".json", ".xml", ".yaml", ".yml", ".h", ".hh", ".hpp", ".inc", ".snippet", ".snippets", ".asm", ".s", ".se", ".sym", ".ini", ".inf", ".map", ".bat", ".rtf"]    
    
    @staticmethod
    def read_pcap_file(file_path):
        import dpkt
        result = ""  # Create an empty string to store the packet details
        with open(file_path, 'rb') as f:
            pcap = dpkt.pcap.Reader(f)
            for timestamp, buf in pcap:
                eth = dpkt.ethernet.Ethernet(buf)
                
                # Extract Ethernet information
                src_mac = ':'.join('{:02x}'.format(b) for b in eth.src)
                dst_mac = ':'.join('{:02x}'.format(b) for b in eth.dst)
                eth_type = eth.type
                
                # Concatenate Ethernet information to the result string
                result += f"Timestamp: {timestamp}\n"
                result += f"Source MAC: {src_mac}\n"
                result += f"Destination MAC: {dst_mac}\n"
                result += f"Ethernet Type: {eth_type}\n"
                
                # Check if packet is an IP packet
                if isinstance(eth.data, dpkt.ip.IP):
                    ip = eth.data
                    
                    # Extract IP information
                    src_ip = dpkt.ip.inet_to_str(ip.src)
                    dst_ip = dpkt.ip.inet_to_str(ip.dst)
                    ip_proto = ip.p
                
                    # Concatenate IP information to the result string
                    result += f"Source IP: {src_ip}\n"
                    result += f"Destination IP: {dst_ip}\n"
                    result += f"IP Protocol: {ip_proto}\n"
                    
                    # Check if packet is a TCP packet
                    if isinstance(ip.data, dpkt.tcp.TCP):
                        tcp = ip.data
                        
                        # Extract TCP information
                        src_port = tcp.sport
                        dst_port = tcp.dport
                        
                        # Concatenate TCP information to the result string
                        result += f"Source Port: {src_port}\n"
                        result += f"Destination Port: {dst_port}\n"
                        
                        # Add more code here to extract and concatenate other TCP details if needed
                        
                    # Check if packet is a UDP packet
                    elif isinstance(ip.data, dpkt.udp.UDP):
                        udp = ip.data
                        
                        # Extract UDP information
                        src_port = udp.sport
                        dst_port = udp.dport
                        
                        # Concatenate UDP information to the result string
                        result += f"Source Port: {src_port}\n"
                        result += f"Destination Port: {dst_port}\n"
                        
                        # Add more code here to extract and concatenate other UDP details if needed
                        
                    # Add more code here to handle other protocols if needed
                    
                result += '-' * 50 + '\n'  # Separator between packets
                
        return result  # Return the result string

    @staticmethod
    def read_pdf_file(file_path: Path) -> str:
        """
        Read a PDF file and return its content as a string.

        Args:
            file_path (Path): The path to the PDF file.

        Returns:
            str: The content of the PDF file.
        """
        import PyPDF2
        def extract_text_from_pdf(file_path):
            text = ""
            with open(file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text += page.extract_text()
            return text
        # Extract text from the PDF
        text = extract_text_from_pdf(file_path)

        # Convert to Markdown (You may need to implement custom logic based on your specific use case)
        markdown_text = text.replace('\n', '  \n')  # Adding double spaces at the end of each line for Markdown line breaks
        
        return markdown_text

    @staticmethod
    def read_docx_file(file_path: Path) -> str:
        """
        Read a DOCX file and return its content as a string.

        Args:
            file_path (Path): The path to the DOCX file.

        Returns:
            str: The content of the DOCX file.
        """
        try:
            from docx import Document
        except ImportError:
            PackageManager.install_package("python-docx")
            from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path: Path) -> str:
        """
        Read a JSON file and return its content as a string.

        Args:
            file_path (Path): The path to the JSON file.

        Returns:
            str: The content of the JSON file.
        """
        import json
        with open(file_path, 'r', encoding='utf-8') as file:
            data = str(json.load(file))
        return data
    
    @staticmethod
    def read_csv_file(file_path: Path) -> str:
        """
        Read a CSV file and return its content as a string.
        Args:
            file_path (Path): The path to the CSV file.
        Returns:
            str: The content of the CSV file.
        """
        with open(file_path, 'r') as file:
            content = file.read()
        return content   

    @staticmethod
    def read_html_file(file_path: Path) -> str:
        """
        Read an HTML file and return its content as a string.

        Args:
            file_path (Path): The path to the HTML file.

        Returns:
            str: The content of the HTML file.
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            PackageManager.install_package("beautifulsoup4")
            from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    
    @staticmethod
    def read_pptx_file(file_path: Path) -> str:
        """
        Read a PPTX file and return its content as a string.

        Args:
            file_path (Path): The path to the PPTX file.

        Returns:
            str: The content of the PPTX file.
        """
        try:
            from pptx import Presentation
        except ImportError:
            PackageManager.install_package("python-pptx")
            from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    
    @staticmethod
    def read_text_file(file_path: Path) -> str:
        """
        Read a text file and return its content as a string.

        Args:
            file_path (Path): The path to the text file.

        Returns:
            str: The content of the text file.
        """
        # Implementation details omitted for brevity
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content

