# safe_store/indexing/parser.py
from pathlib import Path
from typing import Callable, Dict, Union, List, Optional
from ascii_colors import ASCIIColors
import os
import json
import pipmaster as pm
# Import specific custom exceptions
from ..core.exceptions import ParsingError, FileHandlingError, ConfigurationError
# Mimetypes
import mimetypes

# Protocols
from typing import Protocol, runtime_checkable, BinaryIO, TextIO
@runtime_checkable
class PdfReaderProtocol(Protocol): # noqa
    pages: list
    def __init__(self, stream: Union[bytes, Path, str, BinaryIO], strict: bool = True): ...
@runtime_checkable
class PageObjectProtocol(Protocol): # noqa
    def extract_text(self) -> str: ...
@runtime_checkable
class DocumentProtocol(Protocol): # noqa
    paragraphs: list
    def __init__(self, file_path: Union[str, Path, None] = None): ...
@runtime_checkable
class ParagraphProtocol(Protocol): # noqa
    text: str
@runtime_checkable
class BeautifulSoupProtocol(Protocol): # noqa
    def __init__(self, markup: Union[str, bytes], features: str, **kwargs): ...
    def get_text(self, separator: str = "", strip: bool = False) -> str: ...


# --- Parsing Functions ---
def _process_msg_attachment(data: bytes, filename: str, images: dict) -> str:
    # Guess file type based on filename
    mime_type, _ = mimetypes.guess_type(filename)
    file_ext = filename.lower().split(".")[-1] if "." in filename else ""

    # Try decoding as text if likely to be a text attachment
    if mime_type and mime_type.startswith("text") or file_ext in ("txt", "csv", "md"):
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = data.decode("latin1")
            except Exception:
                text = None
        if text and text.strip():
            return f"[Attachment: {filename}]\n{text[:300]}{'...' if len(text)>300 else ''}"
        else:
            return f"[Attachment: {filename} -- undecodable text]"

    # If it's a PDF, Word, or other known binary type, just summarize
    if file_ext in ("pdf", "docx", "xlsx", "pptx", "zip", "eml", "msg") or \
       (mime_type and (mime_type.startswith("application/") or mime_type.startswith("image/"))):
        # Optionally store images for further processing
        if mime_type and mime_type.startswith("image") and images is not None:
            images[filename] = data  # Store image bytes for future use
        return f"[Attachment: {filename} ({mime_type or file_ext}) not parsed as text]"

    # Otherwise, unknown type
    return f"[Attachment: {filename} (unknown type) -- size {len(data)} bytes]"

def parse_msg(file_path: Union[str, Path]) -> str:
    try:
        import extract_msg
        from bs4 import BeautifulSoup  # Correct import
    except ImportError:
        pm.ensure_packages(["extract-msg", "bs4"])  # Use 'bs4' not 'BeautifulSoup'
        import extract_msg
        from bs4 import BeautifulSoup

    try:
        if not Path(file_path).exists():
            raise FileHandlingError(f"File not found: {file_path}")
            
        try:
            msg = extract_msg.Message(str(file_path))
        except (NotImplementedError, Exception) as e:
            # extract_msg often raises generic exceptions for encrypted files or NotImplementedError
            if "encrypted" in str(e).lower():
                raise ParsingError("Content is encrypted") from e
            raise ParsingError(f"Failed to open MSG file: {e}") from e

        images = {}
        header_lines = []
        if getattr(msg, "subject", ""):
            header_lines.append(f"# {getattr(msg, 'subject', '')}")
        meta = []
        sender = getattr(msg, "sender", None) or getattr(msg, "from_", None)
        if sender:
            meta.append(f"From: {sender}")
        if getattr(msg, "to", ""):
            meta.append(f"To: {getattr(msg, 'to', '')}")
        if getattr(msg, "date", ""):
            meta.append(f"Date: {getattr(msg, 'date', '')}")
        if meta:
            header_lines.append("\n".join(meta))
        header = "\n\n".join(header_lines)

        msg_body = (getattr(msg, "body", "") or "").strip()
        if not msg_body and getattr(msg, "htmlBody", None):
            html_body = getattr(msg, "htmlBody", "")
            if html_body:
                msg_body = BeautifulSoup(html_body, "html.parser").get_text()

        attachment_text_parts: List[str] = [header, msg_body]
        if hasattr(msg, "attachments"):
            for att in msg.attachments:
                att_data = getattr(att, "data", b"")
                att_name = getattr(att, "longFilename", None) or getattr(att, "shortFilename", None) or "attachment"
                text_part = _process_msg_attachment(att_data, att_name, images)
                if text_part:
                    attachment_text_parts.append(text_part)

        extracted_text = "\n\n".join([p for p in attachment_text_parts if p and p.strip()])
        
        # Heuristic check: if the library parsed it but it's empty, it might be an issue, 
        # but the empty check in store.py will handle it.
        return extracted_text

    except ParsingError:
        raise
    except Exception as e:
        raise ParsingError(f"Error processing MSG file: {e}") from e
    
def parse_txt(file_path: Union[str, Path]) -> str:
    """
    Parses a plain text file (UTF-8 encoding).
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse TXT file: {_file_path}")
    try:
        with open(_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        ASCIIColors.debug(f"Successfully parsed TXT file: {_file_path}")
        return content
    except FileNotFoundError as e:
        msg = f"File not found: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except UnicodeDecodeError as e:
        msg = f"Encoding error parsing TXT file {_file_path} as UTF-8: {e}"
        ASCIIColors.error(msg)
        raise ParsingError(msg) from e
    except OSError as e:
        msg = f"OS error reading TXT file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg) from e
    except Exception as e:
        msg = f"Unexpected error parsing TXT file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_json(file_path: Union[str, Path]) -> str:
    """
    Parses a JSON file and pretty-prints it to ensure structural integrity for chunking.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse JSON file: {_file_path}")
    try:
        with open(_file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        # Pretty print with indentation to create lines/blocks
        content = json.dumps(data, indent=2, ensure_ascii=False)
        ASCIIColors.debug(f"Successfully parsed and formatted JSON file: {_file_path}")
        return content
    except json.JSONDecodeError as e:
        msg = f"Invalid JSON in file {_file_path}: {e}"
        ASCIIColors.error(msg)
        raise ParsingError(msg) from e
    except Exception as e:
        msg = f"Error parsing JSON file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_csv(file_path: Union[str, Path]) -> str:
    """
    Parses a CSV file intelligently, preserving header context for each row.
    Format: Row {i}: {Header1}={Value1} | {Header2}={Value2} ...
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse CSV file: {_file_path}")
    import csv
    
    try:
        lines = []
        with open(_file_path, 'r', encoding='utf-8', newline='') as f:
            # Attempt to determine dialect
            try:
                sample = f.read(1024)
                f.seek(0)
                dialect = csv.Sniffer().sniff(sample)
                has_header = csv.Sniffer().has_header(sample)
            except Exception:
                # Fallback defaults
                f.seek(0)
                dialect = 'excel'
                has_header = True

            reader = csv.reader(f, dialect=dialect)
            rows = list(reader)

            if not rows:
                return ""

            if has_header:
                headers = rows[0]
                data_rows = rows[1:]
            else:
                headers = [f"Col_{i}" for i in range(len(rows[0]))]
                data_rows = rows

            for i, row in enumerate(data_rows):
                # Pair headers with values
                row_parts = []
                for h_idx, val in enumerate(row):
                    header = headers[h_idx] if h_idx < len(headers) else f"Col_{h_idx}"
                    if val.strip(): # Only add non-empty values
                        row_parts.append(f"{header}={val.strip()}")
                
                if row_parts:
                    lines.append(f"Row {i+1}: " + " | ".join(row_parts))

        ASCIIColors.debug(f"Successfully parsed CSV file: {_file_path}")
        return "\n".join(lines)
    except Exception as e:
        msg = f"Error parsing CSV file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_excel(file_path: Union[str, Path]) -> str:
    """
    Parses an Excel file (XLSX/XLS) intelligently, preserving header context.
    Uses pandas if available.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse Excel file: {_file_path}")
    try:
        pm.ensure_packages(["pandas", "openpyxl"])
        import pandas as pd
    except ImportError as e:
        raise ConfigurationError("Parsing Excel files requires 'pandas' and 'openpyxl'.") from e

    try:
        # Read all sheets
        xls = pd.ExcelFile(_file_path)
        all_text = []

        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            sheet_lines = [f"--- Sheet: {sheet_name} ---"]
            
            # Use columns as headers
            headers = [str(c) for c in df.columns]
            
            for idx, row in df.iterrows():
                row_parts = []
                for h, val in zip(headers, row):
                    if pd.notna(val) and str(val).strip():
                        row_parts.append(f"{h}={str(val).strip()}")
                
                if row_parts:
                    # Using idx+2 because Excel rows start at 1 and header is usually row 1
                    sheet_lines.append(f"Row {idx+2}: " + " | ".join(row_parts))
            
            all_text.append("\n".join(sheet_lines))

        ASCIIColors.debug(f"Successfully parsed Excel file: {_file_path}")
        return "\n\n".join(all_text)
    except Exception as e:
        msg = f"Error parsing Excel file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e


def parse_pdf(file_path: Union[str, Path]) -> str:
    """
    Parses a PDF file to extract text content using pypdf.
    """
    _file_path = Path(file_path)
    ASCIIColors.debug(f"Attempting to parse PDF file: {_file_path}")
    try:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError, EmptyFileError
    except ImportError as e:
        msg = "Parsing PDF files requires 'pypdf'. Install with: pip install safe_store[parsing]"
        ASCIIColors.error(msg)
        raise ConfigurationError(msg) from e

    full_text = ""
    try:
        reader: PdfReaderProtocol = PdfReader(_file_path, strict=False)
        num_pages = len(reader.pages)
        if num_pages == 0:
             return ""

        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    full_text += page_text + "\n"
            except Exception:
                 pass

        return full_text.strip()
    except Exception as e:
        msg = f"Error parsing PDF file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_pptx(file_path: Union[str, Path]) -> str:
    """
    Parses a PPTX file to extract text content using python-pptx.
    """
    _file_path = Path(file_path)
    try:
        from pptx import Presentation
    except ImportError:
        pm.ensure_packages("python-pptx")
        from pptx import Presentation

    full_text = ""
    try:
        presentation = Presentation(_file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"
        return full_text.strip()
    except Exception as e:
        msg = f"Error parsing PPTX file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_docx(file_path: Union[str, Path]) -> str:
    """
    Parses a DOCX file to extract text content using python-docx.
    """
    _file_path = Path(file_path)
    try:
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError:
        pm.ensure_packages("python-docx")
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError

    full_text = ""
    try:
        document: DocumentProtocol = Document(_file_path)
        for para in document.paragraphs:
            full_text += para.text + "\n"
        return full_text.strip()
    except Exception as e:
        msg = f"Error parsing DOCX file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e

def parse_html(file_path: Union[str, Path]) -> str:
    """
    Parses an HTML file to extract text content using BeautifulSoup.
    """
    _file_path = Path(file_path)
    try:
        from bs4 import BeautifulSoup
        try:
            import lxml # noqa F401
            HTML_PARSER = 'lxml'
        except ImportError:
            HTML_PARSER = 'html.parser'
    except ImportError as e:
        msg = "Parsing HTML files requires 'BeautifulSoup4'. Install with: pip install safe_store[parsing]"
        ASCIIColors.error(msg)
        raise ConfigurationError(msg) from e

    try:
        with open(_file_path, 'r', encoding='utf-8') as f:
             content = f.read()
        soup: BeautifulSoupProtocol = BeautifulSoup(content, HTML_PARSER)
        text = soup.get_text(separator='\n', strip=True)
        return text
    except Exception as e:
        msg = f"Error parsing HTML file {_file_path}: {e}"
        ASCIIColors.error(msg, exc_info=True)
        raise ParsingError(msg) from e


# --- Dispatcher Function ---
ParserFunc = Callable[[Union[str, Path]], str]

# Map extensions to parser functions
parser_map: Dict[str, ParserFunc] = {
    '.txt': parse_txt,
    '.pdf': parse_pdf,
    '.docx': parse_docx,
    '.html': parse_html,
    '.htm': parse_html,
    '.pptx': parse_pptx,
    '.msg': parse_msg,

    # Data & Structured
    '.csv': parse_csv,
    '.json': parse_json,
    '.xlsx': parse_excel,
    '.xls': parse_excel,

    # --- General Text & Document Formats ---
    '.md': parse_txt,
    '.rst': parse_txt,
    '.tex': parse_txt,
    '.rtf': parse_txt,
    '.log': parse_txt,
    '.text': parse_txt,
    '.me': parse_txt,
    '.org': parse_txt,

    # --- Configs ---
    '.xml': parse_txt,
    '.tsv': parse_txt, # Should eventually use parse_csv logic with delimiter
    '.yaml': parse_txt,
    '.yml': parse_txt,
    '.sql': parse_txt,
    '.graphql': parse_txt,
    '.gql': parse_txt,
    '.ini': parse_txt,
    '.cfg': parse_txt,
    '.conf': parse_txt,
    '.toml': parse_txt,
    '.env': parse_txt,
    '.properties': parse_txt,
    '.dockerfile': parse_txt,
    'dockerfile': parse_txt,
    '.gitattributes': parse_txt,
    '.gitconfig': parse_txt,
    '.gitignore': parse_txt,

    # --- Programming Languages (Source is text) ---
    '.py': parse_txt, '.pyw': parse_txt, '.js': parse_txt, '.mjs': parse_txt,
    '.cjs': parse_txt, '.ts': parse_txt, '.tsx': parse_txt, '.rb': parse_txt,
    '.pl': parse_txt, '.pm': parse_txt, '.php': parse_txt, '.phtml': parse_txt,
    '.sh': parse_txt, '.bash': parse_txt, '.zsh': parse_txt, '.fish': parse_txt,
    '.ps1': parse_txt, '.psm1': parse_txt, '.psd1': parse_txt, '.lua': parse_txt,
    '.r': parse_txt, '.tcl': parse_txt, '.dart': parse_txt, '.c': parse_txt,
    '.h': parse_txt, '.cpp': parse_txt, '.cxx': parse_txt, '.cc': parse_txt,
    '.hpp': parse_txt, '.hxx': parse_txt, '.cs': parse_txt, '.java': parse_txt,
    '.go': parse_txt, '.rs': parse_txt, '.swift': parse_txt, '.kt': parse_txt,
    '.kts': parse_txt, '.scala': parse_txt, '.m': parse_txt, '.mm': parse_txt,
    '.f': parse_txt, '.for': parse_txt, '.f90': parse_txt, '.f95': parse_txt,
    '.pas': parse_txt, '.d': parse_txt, '.vb': parse_txt, '.vbs': parse_txt,
    '.css': parse_txt, '.scss': parse_txt, '.sass': parse_txt, '.less': parse_txt,
    '.styl': parse_txt, '.jsx': parse_txt, '.vue': parse_txt, '.svelte': parse_txt,
    '.ejs': parse_txt, '.hbs': parse_txt, '.mustache': parse_txt, '.jinja': parse_txt,
    '.jinja2': parse_txt, '.twig': parse_txt, '.erb': parse_txt, '.jsp': parse_txt,
    '.asp': parse_txt, '.aspx': parse_txt, '.asm': parse_txt, '.s': parse_txt,
    '.bat': parse_txt, '.cmd': parse_txt, '.vhd': parse_txt, '.vhdl': parse_txt,
    '.sv': parse_txt, '.bib': parse_txt, '.srt': parse_txt, '.sub': parse_txt,
    '.vtt': parse_txt, '.po': parse_txt, '.pot': parse_txt, '.strings': parse_txt,
}

SAFE_STORE_SUPPORTED_FILE_EXTENSIONS = list(parser_map.keys())

def parse_document(file_path: Union[str, Path]) -> str:
    """
    Parses a document based on its file extension.
    """
    _file_path = Path(file_path)
    if not _file_path.is_file():
        msg = f"Input path is not a file: {_file_path}"
        ASCIIColors.error(msg)
        raise FileHandlingError(msg)

    extension = _file_path.suffix.lower()
    ASCIIColors.debug(f"Dispatching parser for extension '{extension}' on file: {_file_path.name}")

    parser_func = parser_map.get(extension)

    if parser_func:
        try:
            return parser_func(_file_path)
        except (ConfigurationError, FileHandlingError, ParsingError) as e:
             raise e
        except Exception as e:
             msg = f"Unexpected error during parsing dispatch for {_file_path.name}: {e}"
             ASCIIColors.error(msg, exc_info=True)
             raise ParsingError(msg) from e
    else:
        msg = f"Unsupported file type extension: '{extension}' for file: {_file_path}."
        ASCIIColors.warning(msg)
        raise ConfigurationError(msg)
